#!/usr/bin/env python3
"""
Build Week11_AI_Trends_Slides.pptx from the week11-slides-batch*.md files.

Content rules (prevent overflow):
  - Each slide shows ONE primary section: either the first table (≤7 rows) OR
    a curated set of bullets (≤6), never both large blocks together.
  - Chart slides show abbreviated content in the left pane only.
  - Height budget is tracked precisely; content stops when space runs out.
  - Code blocks limited to 14 lines.

Usage:
    pip install python-pptx
    python3 build_pptx.py

Output:
    Week11_AI_Trends_Slides.pptx
"""

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Chart slides (have a charts/slide_NN_chart.png) ───────────────────────────
CHART_SLIDES = {3, 4, 5, 6, 7, 14, 20, 25, 28, 29, 30}

# ── Course colors ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x16, 0x21, 0x3E)
NAVY_MID    = RGBColor(0x0F, 0x34, 0x60)
PURPLE      = RGBColor(0x53, 0x34, 0x83)
TEAL        = RGBColor(0x0E, 0xA5, 0xE9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF8, 0xF9, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
MID_TEXT    = RGBColor(0x2D, 0x37, 0x48)
CODE_BG     = RGBColor(0x1A, 0x1A, 0x2E)
CODE_TEXT   = RGBColor(0xE2, 0xE8, 0xF0)
AMBER_RGB   = RGBColor(0xF5, 0x9E, 0x0B)

# ── Slide geometry (16:9, 13.333 × 7.5 inches) ────────────────────────────────
SLIDE_W      = Inches(13.333)
SLIDE_H      = Inches(7.5)
TITLE_H      = Inches(1.35)
CONTENT_TOP  = Inches(1.48)
CONTENT_H    = Inches(5.75)
MARGIN_L     = Inches(0.55)
MARGIN_R     = Inches(0.55)
CONTENT_W    = SLIDE_W - MARGIN_L - MARGIN_R
SLIDE_BOTTOM = SLIDE_H - Inches(0.45)   # hard lower boundary

# ── Height budget per block type (in Emu) ─────────────────────────────────────
H_H3        = Inches(0.44)
H_H4        = Inches(0.34)
H_BULLET    = Inches(0.32)
H_SUBBULLET = Inches(0.27)
H_PARA      = Inches(0.30)
H_QUOTE     = Inches(0.62)
H_GAP       = Inches(0.10)
H_TABLE_ROW = Inches(0.36)
H_CODE_LINE = Inches(0.265)
H_CODE_PAD  = Inches(0.28)   # top+bottom padding in code box

MAX_TABLE_ROWS_CHART = 5    # max rows shown in chart-slide left pane
MAX_TABLE_ROWS_FULL  = 7    # max rows shown in full-width slide
MAX_BULLETS_CHART    = 5    # max bullet items in chart-slide left pane
MAX_CODE_LINES       = 14


# ── Data structures ────────────────────────────────────────────────────────────

@dataclass
class SlideContent:
    title: str = ""
    subtitle: str = ""
    blocks: list = field(default_factory=list)


@dataclass
class TextBlock:
    kind: str    # "para", "bullet", "subbullet", "h3", "h4"
    text: str


@dataclass
class CodeBlock:
    kind: str = "code"
    language: str = ""
    lines: list = field(default_factory=list)


@dataclass
class QuoteBlock:
    kind: str = "quote"
    text: str = ""


@dataclass
class TableBlock:
    kind: str = "table"
    headers: list = field(default_factory=list)
    rows: list = field(default_factory=list)


# ── Markdown parser ────────────────────────────────────────────────────────────

def strip_inline(text: str) -> str:
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*',     r'\1', text)
    text = re.sub(r'`(.+?)`',       r'\1', text)
    return text.strip()


def parse_batches(batch_files: list) -> list:
    slides = []
    current = None
    in_code = False
    code_lang = ""
    code_lines = []
    in_table = False
    table_headers = []
    table_rows = []

    def flush_table():
        nonlocal in_table, table_headers, table_rows
        if in_table and current is not None:
            current.blocks.append(TableBlock(headers=table_headers, rows=table_rows))
        in_table = False
        table_headers = []
        table_rows = []

    def flush_code():
        nonlocal in_code, code_lang, code_lines
        if in_code and current is not None:
            current.blocks.append(CodeBlock(language=code_lang, lines=list(code_lines)))
        in_code = False
        code_lang = ""
        code_lines = []

    for path in batch_files:
        if not path.exists():
            print(f"  ✗ Missing: {path.name}")
            continue
        for raw_line in path.read_text(encoding="utf-8").splitlines():
            line = raw_line.rstrip()

            if line.startswith("```"):
                if not in_code:
                    flush_table()
                    in_code = True
                    code_lang = line[3:].strip() or "text"
                    code_lines = []
                else:
                    flush_code()
                continue

            if in_code:
                code_lines.append(line)
                continue

            if "|" in line and line.strip().startswith("|"):
                if re.match(r"^[\|\s\-:]+$", line):
                    continue
                cells = [strip_inline(c) for c in line.strip().strip("|").split("|")]
                if not in_table:
                    flush_table()
                    in_table = True
                    table_headers = cells
                else:
                    table_rows.append(cells)
                continue
            else:
                if in_table:
                    flush_table()

            if line.startswith("## Slide "):
                flush_code(); flush_table()
                raw_title = line[3:].strip()
                m = re.match(r"Slide\s+\d+:\s+(.*)", raw_title)
                clean_title = m.group(1) if m else raw_title
                current = SlideContent(title=clean_title)
                slides.append(current)
                continue

            if line.startswith("# ") or line.strip() == "---":
                continue
            if current is None:
                continue

            if line.startswith("### "):
                content = strip_inline(line[4:])
                if not current.subtitle:
                    current.subtitle = content
                else:
                    current.blocks.append(TextBlock(kind="h3", text=content))
                continue

            if line.startswith("#### "):
                current.blocks.append(TextBlock(kind="h4", text=strip_inline(line[5:])))
                continue

            if line.startswith("> "):
                current.blocks.append(QuoteBlock(text=strip_inline(line[2:])))
                continue

            m = re.match(r"^([ \t]*)[-\*] (.+)$", line)
            if m:
                indent = len(m.group(1).expandtabs(4))
                kind = "subbullet" if indent >= 2 else "bullet"
                current.blocks.append(TextBlock(kind=kind, text=strip_inline(m.group(2))))
                continue

            m2 = re.match(r"^\d+\. (.+)$", line)
            if m2:
                current.blocks.append(TextBlock(kind="bullet",
                                                 text=strip_inline(m2.group(1))))
                continue

            stripped = line.strip()
            if stripped:
                current.blocks.append(TextBlock(kind="para", text=strip_inline(stripped)))

    flush_code(); flush_table()
    return slides


# ── Content selection ──────────────────────────────────────────────────────────

def _curate(blocks: list, has_chart: bool) -> list:
    """
    Return a curated subset of blocks that will fit on one slide.

    Strategy:
      - Show subtitle (h3) if present, skip duplicates
      - Show at most ONE table (trimmed to row limit)
      - Show at most ONE code block (trimmed to line limit)
      - Show limited bullets/paras
      - Never mix a large table + code block on the same slide
    """
    max_tbl = MAX_TABLE_ROWS_CHART if has_chart else MAX_TABLE_ROWS_FULL
    max_blt = MAX_BULLETS_CHART    if has_chart else 7

    out = []
    n_tables = 0
    n_codes  = 0
    n_bullets = 0
    seen_h3 = set()

    for blk in blocks:
        if isinstance(blk, TableBlock):
            if n_tables >= 1:
                continue
            # Trim rows
            trimmed = TableBlock(headers=blk.headers,
                                  rows=blk.rows[:max_tbl])
            out.append(trimmed)
            n_tables += 1

        elif isinstance(blk, CodeBlock):
            if n_codes >= 1 or n_tables >= 1:
                continue
            trimmed = CodeBlock(language=blk.language,
                                 lines=blk.lines[:MAX_CODE_LINES])
            out.append(trimmed)
            n_codes += 1

        elif isinstance(blk, QuoteBlock):
            if n_bullets < max_blt:
                out.append(blk)

        elif isinstance(blk, TextBlock):
            if blk.kind == "h3":
                if blk.text not in seen_h3:
                    seen_h3.add(blk.text)
                    out.append(blk)
            elif blk.kind == "h4":
                out.append(blk)
            elif blk.kind in ("bullet", "subbullet"):
                if n_bullets < max_blt:
                    out.append(blk)
                    n_bullets += 1
            elif blk.kind == "para":
                if n_bullets < max_blt:
                    out.append(blk)
                    n_bullets += 1

    return out


# ── python-pptx helpers ────────────────────────────────────────────────────────

def solid_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_title_bar(slide, title_text: str, subtitle_text: str = ""):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, TITLE_H)
    solid_fill(bar, NAVY)
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.45), Inches(0.1), SLIDE_W - Inches(0.9), Inches(0.72))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.color.rgb = WHITE
    run.font.size = Pt(27)
    run.font.bold = True
    run.font.name = "Calibri"

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(
            Inches(0.45), Inches(0.82), SLIDE_W - Inches(0.9), Inches(0.44))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle_text
        srun.font.color.rgb = TEAL
        srun.font.size = Pt(15)
        srun.font.name = "Calibri"

    line = slide.shapes.add_shape(1, Inches(0), TITLE_H, SLIDE_W, Inches(0.04))
    solid_fill(line, TEAL)
    line.line.fill.background()


def add_bg(slide):
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    solid_fill(bg, WHITE)
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_slide_num(slide, num: int):
    nb = slide.shapes.add_textbox(
        SLIDE_W - Inches(0.7), SLIDE_H - Inches(0.35),
        Inches(0.55), Inches(0.28))
    np_ = nb.text_frame.paragraphs[0]
    np_.alignment = PP_ALIGN.CENTER
    nr = np_.add_run()
    nr.text = str(num)
    nr.font.size = Pt(10)
    nr.font.color.rgb = RGBColor(0xAA, 0xAA, 0xBB)
    nr.font.name = "Calibri"


def render_table(slide, tbl: TableBlock, left, top, width, max_h):
    n_cols = max(len(tbl.headers),
                 max((len(r) for r in tbl.rows), default=0))
    n_rows = 1 + len(tbl.rows)
    if n_cols == 0:
        return Inches(0)

    row_h = min(Inches(0.36), max_h // n_rows)
    tbl_h = row_h * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, tbl_h).table

    for ci, hdr in enumerate(tbl.headers[:n_cols]):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.color.rgb = WHITE
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.name = "Calibri"

    for ri, row in enumerate(tbl.rows, start=1):
        row_color = LIGHT_GRAY if ri % 2 == 0 else WHITE
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.text = row[ci] if ci < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            para = cell.text_frame.paragraphs[0]
            for run in para.runs:
                run.font.size = Pt(10)
                run.font.name = "Calibri"
                run.font.color.rgb = DARK_TEXT

    return tbl_h


def render_code(slide, blk: CodeBlock, left, top, width):
    lines = blk.lines[:MAX_CODE_LINES]
    n_lines = max(len(lines), 1)
    box_h = H_CODE_LINE * n_lines + H_CODE_PAD

    box = slide.shapes.add_shape(1, left, top, width, box_h)
    solid_fill(box, CODE_BG)
    box.line.fill.background()

    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), box_h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False

    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        run.font.color.rgb = CODE_TEXT
        run.font.size = Pt(9)
        run.font.name = "Courier New"

    return box_h


# ── Core block renderer ────────────────────────────────────────────────────────

def _render_blocks(slide, blocks: list, left, top, width):
    """
    Render a curated list of blocks into the slide, tracking height budget.
    Returns the final `top` position after all content.
    """
    BOTTOM = SLIDE_BOTTOM
    content_tf_box = None
    content_tf = None

    def ensure_text_box():
        nonlocal content_tf_box, content_tf
        if content_tf_box is None:
            remaining = max(BOTTOM - top, Inches(0.5))
            content_tf_box = slide.shapes.add_textbox(left, top, width, remaining)
            content_tf = content_tf_box.text_frame
            content_tf.word_wrap = True

    def flush_text_box():
        nonlocal content_tf_box, content_tf, top
        if content_tf_box is not None:
            n = len(content_tf.paragraphs)
            est = H_BULLET * n + H_GAP
            top = min(top + est, BOTTOM)
            content_tf_box = None
            content_tf = None

    def add_para(text, size, bold=False, color=None, indent=0, bullet_char=""):
        ensure_text_box()
        paras = content_tf.paragraphs
        p = paras[0] if (len(paras) == 1 and not paras[0].runs) else \
            content_tf.add_paragraph()
        p.level = indent
        p.space_before = Pt(3)
        p.space_after  = Pt(2)
        if bullet_char:
            rb = p.add_run()
            rb.text = bullet_char + " "
            rb.font.size = Pt(size)
            rb.font.color.rgb = TEAL if bullet_char in ("•", "›") else NAVY_MID
            rb.font.name = "Calibri"
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = color or DARK_TEXT

    for blk in blocks:
        if top >= BOTTOM - Inches(0.3):
            break

        if isinstance(blk, TableBlock):
            flush_text_box()
            avail = BOTTOM - top
            if avail < Inches(0.5):
                break
            tbl_h = render_table(slide, blk, left, top, width, avail)
            top += tbl_h + H_GAP

        elif isinstance(blk, CodeBlock):
            flush_text_box()
            needed = H_CODE_LINE * min(len(blk.lines), MAX_CODE_LINES) + H_CODE_PAD
            if top + needed > BOTTOM:
                break
            used = render_code(slide, blk, left, top, width)
            top += used + H_GAP

        elif isinstance(blk, QuoteBlock):
            flush_text_box()
            if top + H_QUOTE > BOTTOM:
                break
            bar = slide.shapes.add_shape(1, left, top, Inches(0.06), H_QUOTE * 0.88)
            solid_fill(bar, PURPLE)
            bar.line.fill.background()
            qtb = slide.shapes.add_textbox(
                left + Inches(0.18), top,
                width - Inches(0.18), H_QUOTE)
            qtf = qtb.text_frame
            qtf.word_wrap = True
            qp = qtf.paragraphs[0]
            qrun = qp.add_run()
            qrun.text = blk.text
            qrun.font.size = Pt(13)
            qrun.font.italic = True
            qrun.font.color.rgb = NAVY_MID
            qrun.font.name = "Calibri"
            top += H_QUOTE + H_GAP

        elif isinstance(blk, TextBlock):
            if blk.kind == "h3":
                flush_text_box()
                if top + H_H3 > BOTTOM:
                    break
                add_para(blk.text, 16, bold=True, color=PURPLE)
            elif blk.kind == "h4":
                if top + H_H4 > BOTTOM:
                    break
                add_para(blk.text, 13, bold=True, color=NAVY_MID)
            elif blk.kind == "bullet":
                if top + H_BULLET > BOTTOM:
                    break
                add_para(blk.text, 13, bullet_char="•")
            elif blk.kind == "subbullet":
                if top + H_SUBBULLET > BOTTOM:
                    break
                add_para(blk.text, 11, bullet_char="›", indent=1, color=MID_TEXT)
            elif blk.kind == "para":
                if top + H_PARA > BOTTOM:
                    break
                if len(blk.text) < 90 and blk.text.endswith(":"):
                    add_para(blk.text, 13, bold=True, color=NAVY_MID)
                else:
                    add_para(blk.text, 12, color=MID_TEXT)

    flush_text_box()
    return top


# ── Slide renderers ────────────────────────────────────────────────────────────

# Split-layout geometry
TEXT_PANE_W  = Inches(7.0)
CHART_PANE_L = Inches(7.7)
CHART_PANE_W = Inches(5.1)
DIVIDER_X    = Inches(7.6)


def _chart_path(script_dir: Path, slide_num: int):
    p = script_dir / "charts" / f"slide_{slide_num:02d}_chart.png"
    return p if p.exists() else None


def render_slide_with_chart(prs, sc: SlideContent, slide_num: int, chart_img: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, sc.title, sc.subtitle)
    add_slide_num(slide, slide_num)

    # Divider
    div = slide.shapes.add_shape(1, DIVIDER_X, CONTENT_TOP, Inches(0.02), CONTENT_H)
    solid_fill(div, RGBColor(0xDD, 0xDD, 0xEE))
    div.line.fill.background()

    # Chart image — right pane
    slide.shapes.add_picture(
        str(chart_img),
        CHART_PANE_L, CONTENT_TOP + Inches(0.1),
        CHART_PANE_W, CONTENT_H - Inches(0.2))

    # Left pane — curated text
    curated = _curate(sc.blocks, has_chart=True)
    _render_blocks(slide, curated, MARGIN_L, CONTENT_TOP, TEXT_PANE_W)


def render_slide(prs, sc: SlideContent, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, sc.title, sc.subtitle)
    add_slide_num(slide, slide_num)

    curated = _curate(sc.blocks, has_chart=False)
    _render_blocks(slide, curated, MARGIN_L, CONTENT_TOP, CONTENT_W)


# ── Cover slide ────────────────────────────────────────────────────────────────

def render_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg1 = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    solid_fill(bg1, NAVY)
    bg1.line.fill.background()

    accent = slide.shapes.add_shape(1, Inches(0), Inches(4.5), SLIDE_W, Inches(3.0))
    solid_fill(accent, NAVY_MID)
    accent.line.fill.background()

    tbar = slide.shapes.add_shape(1, Inches(0), Inches(3.8), SLIDE_W, Inches(0.08))
    solid_fill(tbar, TEAL)
    tbar.line.fill.background()

    ttb = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.8), Inches(1.6))
    tp = ttb.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = "Week 11: New AI Tools and Trends"
    tr.font.size = Pt(42)
    tr.font.bold = True
    tr.font.color.rgb = WHITE
    tr.font.name = "Calibri"

    stb = slide.shapes.add_textbox(Inches(0.8), Inches(2.95), Inches(11.8), Inches(0.8))
    sp = stb.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "The Full 2026 AI Landscape — Models · Agents · Coding Tools · Automation · Strategy"
    sr.font.size = Pt(19)
    sr.font.color.rgb = TEAL
    sr.font.name = "Calibri"

    mtb = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.8), Inches(1.8))
    mtf = mtb.text_frame
    for i, (line, size, color) in enumerate([
        ("BUAN 6v99 — Generative AI for Business", 16, WHITE),
        ("University of Texas at Dallas  |  Spring 2026", 14, RGBColor(0xAA, 0xBB, 0xDD)),
        ("April 8, 2026  |  Professor Antonio de Pádua Paes Jr.", 14, RGBColor(0xAA, 0xBB, 0xDD)),
    ]):
        mp = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        mr = mp.add_run()
        mr.text = line
        mr.font.size = Pt(size)
        mr.font.color.rgb = color
        mr.font.name = "Calibri"
        mp.space_before = Pt(5)


# ── Section divider ────────────────────────────────────────────────────────────

SECTION_TITLES = {
    7:  ("Section 2", "AI Agent Frameworks & MCP"),
    14: ("Section 3", "AI Coding Tools"),
    20: ("Section 4", "Business Automation & Multimodal AI"),
    26: ("Section 5", "Key Trends & Strategic Outlook"),
}


def render_section_divider(prs, section_num: str, section_title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    solid_fill(bg, NAVY_MID)
    bg.line.fill.background()

    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08))
    solid_fill(bar, TEAL)
    bar.line.fill.background()

    nb = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.0), Inches(0.7))
    nr = nb.text_frame.paragraphs[0].add_run()
    nr.text = section_num
    nr.font.size = Pt(22)
    nr.font.color.rgb = TEAL
    nr.font.name = "Calibri"
    nr.font.bold = True

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.85), Inches(11.0), Inches(1.2))
    trun = tb.text_frame.paragraphs[0].add_run()
    trun.text = section_title
    trun.font.size = Pt(36)
    trun.font.bold = True
    trun.font.color.rgb = WHITE
    trun.font.name = "Calibri"


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    script_dir = Path(__file__).parent
    batch_files = [script_dir / f"week11-slides-batch{i}.md" for i in range(1, 6)]

    print("Parsing markdown batches...")
    slides = parse_batches(batch_files)
    print(f"  Found {len(slides)} slides")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    render_title_slide(prs)
    print("  ✓ Cover slide")

    charts_used = 0
    for i, sc in enumerate(slides, start=1):
        if i in SECTION_TITLES:
            sec_num, sec_title = SECTION_TITLES[i]
            render_section_divider(prs, sec_num, sec_title)
            print(f"  ✓ Section divider: {sec_title}")

        chart_img = _chart_path(script_dir, i)
        if chart_img:
            render_slide_with_chart(prs, sc, i, chart_img)
            print(f"  ✓ Slide {i:2d} [+chart]: {sc.title[:55]}")
            charts_used += 1
        else:
            render_slide(prs, sc, i)
            print(f"  ✓ Slide {i:2d}:         {sc.title[:60]}")

    output = script_dir / "Week11_AI_Trends_Slides.pptx"
    prs.save(str(output))
    size_mb = output.stat().st_size / (1024 * 1024)
    print(f"\n✅  Saved: {output.name}")
    print(f"   Total slides: {len(prs.slides)}  |  Charts embedded: {charts_used}")
    print(f"   File size: {size_mb:.2f} MB")


if __name__ == "__main__":
    main()
