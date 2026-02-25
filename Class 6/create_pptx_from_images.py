#!/usr/bin/env python3
"""
Create PowerPoint presentation from generated slide images.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def create_presentation_from_images(image_dir, output_file):
    """
    Create a PowerPoint presentation from all slide images in a directory.
    
    Args:
        image_dir: Directory containing slide images
        output_file: Output PPTX file path
    """
    # Create a presentation object
    prs = Presentation()
    
    # Set slide width and height (16:9 aspect ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Get all slide images sorted by number
    image_files = sorted([f for f in os.listdir(image_dir) 
                         if f.startswith('slide_') and f.endswith('_final.png')])
    
    print(f"Found {len(image_files)} slide images")
    
    # Add each image as a slide
    for image_file in image_files:
        image_path = os.path.join(image_dir, image_file)
        slide_num = image_file.split('_')[1]
        
        print(f"Adding slide {slide_num}...")
        
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add the image to fill the entire slide
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    
    # Save the presentation
    prs.save(output_file)
    print(f"\n✅ Presentation created successfully: {output_file}")
    print(f"   Total slides: {len(image_files)}")

def main():
    """Main function to create the presentation."""
    # Set paths
    script_dir = Path(__file__).parent
    image_dir = script_dir / "slide_images"
    output_file = script_dir / "Week06_Coding_with_AI_Presentation.pptx"
    
    # Verify image directory exists
    if not image_dir.exists():
        print(f"❌ Error: Image directory not found: {image_dir}")
        return
    
    # Count images
    image_files = list(image_dir.glob("slide_*_final.png"))
    if not image_files:
        print(f"❌ Error: No slide images found in {image_dir}")
        return
    
    print("=" * 60)
    print("Creating PowerPoint Presentation")
    print("=" * 60)
    print(f"Image directory: {image_dir}")
    print(f"Output file: {output_file}")
    print("=" * 60)
    print()
    
    # Create the presentation
    create_presentation_from_images(str(image_dir), str(output_file))
    
    print()
    print("=" * 60)
    print("✨ Done! You can now open the presentation:")
    print(f"   open '{output_file}'")
    print("=" * 60)

if __name__ == "__main__":
    main()
