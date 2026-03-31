#!/usr/bin/env python3
"""
Build Week08_LLM_Slides.pptx from slide images in Class 8/images/.
Each image becomes one full-bleed 16:9 slide.

Usage:
    python3 build_pptx.py

Requires: pip install python-pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

SCRIPT_DIR = Path(__file__).parent
IMAGES_DIR = SCRIPT_DIR / "images"
OUTPUT_FILE = SCRIPT_DIR / "Week08_LLM_Slides.pptx"

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)


def build_pptx():
    images = sorted(IMAGES_DIR.glob("slide_*_final.png"))
    if not images:
        print("No slide images found in images/. Run generate_slide_images.py first.")
        return

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank layout

    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0, top=0,
            width=SLIDE_W, height=SLIDE_H,
        )
        print(f"  ✓ Added {img_path.name}")

    prs.save(OUTPUT_FILE)
    size_mb = OUTPUT_FILE.stat().st_size / (1024 * 1024)
    print(f"\n✅ PPTX saved: {OUTPUT_FILE.name}")
    print(f"   Slides: {len(images)}")
    print(f"   Size:   {size_mb:.1f} MB")


if __name__ == "__main__":
    build_pptx()
